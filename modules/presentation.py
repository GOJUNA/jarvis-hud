import os
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from utils.logger import log


class PresentationManager:
    """Erstellt PowerPoint-Praesentationen programmatisch."""

    OUTPUT_DIR = "data/presentations"

    THEMES = {
        "dark": {
            "bg": RGBColor(0x0A, 0x0E, 0x14),
            "primary": RGBColor(0x00, 0xD4, 0xFF),
            "secondary": RGBColor(0xFF, 0xD7, 0x00),
            "text": RGBColor(0xE6, 0xED, 0xF3),
            "text_dim": RGBColor(0x8B, 0x94, 0x9E),
            "accent": RGBColor(0x00, 0xFF, 0x88),
        },
        "light": {
            "bg": RGBColor(0xFA, 0xFA, 0xFA),
            "primary": RGBColor(0x00, 0x66, 0x99),
            "secondary": RGBColor(0xCC, 0x99, 0x00),
            "text": RGBColor(0x1A, 0x1A, 0x2E),
            "text_dim": RGBColor(0x55, 0x55, 0x55),
            "accent": RGBColor(0x00, 0x88, 0x44),
        },
        "jarvis": {
            "bg": RGBColor(0x03, 0x08, 0x10),
            "primary": RGBColor(0x00, 0xD4, 0xFF),
            "secondary": RGBColor(0xFF, 0xD7, 0x00),
            "text": RGBColor(0xC8, 0xD6, 0xE5),
            "text_dim": RGBColor(0x57, 0x65, 0x74),
            "accent": RGBColor(0x00, 0xFF, 0x88),
        },
    }

    def __init__(self):
        os.makedirs(self.OUTPUT_DIR, exist_ok=True)

    def create_presentation(self, topic: str, slides_count: int = 8, theme: str = "jarvis") -> str:
        """Erstellt eine vollstaendige Praesentation zu einem Thema."""
        try:
            theme_colors = self.THEMES.get(theme, self.THEMES["jarvis"])

            # Hole Recherche-Inhalte
            content = self._generate_content(topic, slides_count)

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Titel-Folie
            self._add_title_slide(prs, content["title"], content["subtitle"], theme_colors)

            # Inhalts-Folien
            for i, slide_data in enumerate(content["slides"]):
                self._add_content_slide(prs, slide_data, theme_colors, i + 1)

            # Abschluss-Folie
            self._add_closing_slide(prs, theme_colors)

            # Speichern
            filename = f"{topic.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = os.path.join(self.OUTPUT_DIR, filename)
            prs.save(filepath)

            log.info(f"Praesentation erstellt: {filepath}")
            return filepath

        except Exception as e:
            log.error(f"Fehler beim Erstellen der Praesentation: {e}")
            raise

    def _generate_content(self, topic: str, slides_count: int) -> dict:
        """Generiert strukturierte Inhalte fuer die Praesentation."""
        # Basis-Struktur basierend auf Thema
        slides = []

        # Intro-Folie
        slides.append({
            "title": f"Was ist {topic}?",
            "bullets": [
                f"Definition und Grundlagen von {topic}",
                "Historische Entwicklung",
                "Warum ist es heute relevant?",
                "Anwendungsbereiche im Ueberblick",
            ],
            "type": "overview",
        })

        # Je nach Thema spezifische Folien
        topic_lower = topic.lower()

        if any(kw in topic_lower for kw in ["programmier", "coding", "software", "entwicklung"]):
            slides.extend(self._programming_slides(topic))
        elif any(kw in topic_lower for kw in ["ki", "künstliche intelligenz", "machine learning", "ml", "ai"]):
            slides.extend(self._ai_slides(topic))
        elif any(kw in topic_lower for kw in ["web", "html", "javascript", "frontend", "backend"]):
            slides.extend(self._web_slides(topic))
        elif any(kw in topic_lower for kw in ["daten", "data", "analyse", "science"]):
            slides.extend(self._data_slides(topic))
        else:
            slides.extend(self._generic_slides(topic))

        # Auf gewuenschte Anzahl kuerzen/erweitern
        slides = slides[:slides_count - 2]  # -2 fuer Titel + Abschluss

        return {
            "title": topic,
            "subtitle": f"Praesentation erstellt von J.A.R.V.I.S. | {datetime.now().strftime('%d.%m.%Y')}",
            "slides": slides,
        }

    def _programming_slides(self, topic: str) -> list:
        return [
            {"title": "Programmierparadigmen", "bullets": ["Imperativ vs. Deklarativ", "Objektorientiert (OOP)", "Funktional", "Prozedural"], "type": "list"},
            {"title": "Wichtige Konzepte", "bullets": ["Variablen & Datentypen", "Kontrollstrukturen (if/for/while)", "Funktionen & Module", "Fehlerbehandlung (try/except)"], "type": "list"},
            {"title": "Moderne Werkzeuge", "bullets": ["Git & Versionskontrolle", "IDEs (VS Code, PyCharm)", "Package Manager (pip, npm)", "CI/CD Pipelines"], "type": "list"},
            {"title": "Best Practices", "bullets": ["Clean Code Prinzipien", "DRY (Don't Repeat Yourself)", "SOLID Prinzipien", "Testing & Documentation"], "type": "list"},
            {"title": "Projekt-Beispiel", "bullets": ["JARVIS - Sprachassistent", "Python + Flask + WebSockets", "Three.js 3D Hologramm", "Web Speech API (STT/TTS)"], "type": "example"},
        ]

    def _ai_slides(self, topic: str) -> list:
        return [
            {"title": "Grundlagen der KI", "bullets": ["Machine Learning vs. Deep Learning", "Neuronale Netze", "Training & Inferenz", "Ueberwachtes vs. unueberwachtes Lernen"], "type": "list"},
            {"title": "Wichtige Architekturen", "bullets": ["CNN (Bilder)", "RNN/LSTM (Sequenzen)", "Transformer (NLP)", "GANs (Generierung)"], "type": "list"},
            {"title": "Anwendungen", "bullets": ["Computer Vision", "NLP & Chatbots", "Empfehlungssysteme", "Autonome Systeme"], "type": "list"},
            {"title": "Herausforderungen", "bullets": ["Bias & Fairness", "Erklaerbarkeit (XAI)", "Rechenaufwand", "Datenschutz"], "type": "list"},
            {"title": "Zukunft", "bullets": ["AGI - Allgemeine KI", "Multimodale Modelle", "KI-Sicherheit", "Mensch-KI-Kollaboration"], "type": "list"},
        ]

    def _web_slides(self, topic: str) -> list:
        return [
            {"title": "Frontend Grundlagen", "bullets": ["HTML5 Semantik", "CSS3 (Flexbox, Grid)", "JavaScript ES6+", "TypeScript"], "type": "list"},
            {"title": "Moderne Frameworks", "bullets": ["React / Vue / Svelte", "Next.js / Nuxt", "State Management", "Component Architecture"], "type": "list"},
            {"title": "Backend & APIs", "bullets": ["REST vs. GraphQL", "Datenbanken (SQL/NoSQL)", "Authentication (JWT/OAuth)", "Microservices"], "type": "list"},
            {"title": "Deployment & DevOps", "bullets": ["Container (Docker)", "Kubernetes", "Serverless", "Monitoring"], "type": "list"},
        ]

    def _data_slides(self, topic: str) -> list:
        return [
            {"title": "Data Pipeline", "bullets": ["Collect -> Clean -> Analyze -> Visualize", "ETL vs. ELT", "Data Quality", "Versioning"], "type": "list"},
            {"title": "Tools & Libraries", "bullets": ["Pandas / NumPy", "Matplotlib / Seaborn / Plotly", "Scikit-learn", "Jupyter Notebooks"], "type": "list"},
            {"title": "Analyse-Methoden", "bullets": ["Explorative Datenanalyse", "Statistische Tests", "Machine Learning", "Time Series"], "type": "list"},
        ]

    def _generic_slides(self, topic: str) -> list:
        return [
            {"title": f"Grundlagen von {topic}", "bullets": ["Kernkonzepte", "Wichtige Begriffe", "Zusammenhaenge", "Grundlagenwissen"], "type": "list"},
            {"title": f"Anwendung von {topic}", "bullets": ["Praxisbeispiele", "Use Cases", "Best Practices", "Tipps & Tricks"], "type": "list"},
            {"title": f"Fortgeschrittenes {topic}", "bullets": ["Expertenwissen", "Spezialthemen", "Aktuelle Trends", "Forschung"], "type": "list"},
        ]

    def _add_title_slide(self, prs, title: str, subtitle: str, colors: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        self._set_slide_bg(slide, colors["bg"])

        # Top accent line
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), colors["primary"])

        # Title
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = colors["text_dim"]
        p2.alignment = PP_ALIGN.CENTER

        # Bottom accent line
        self._add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), colors["primary"])

    def _add_content_slide(self, prs, slide_data: dict, colors: dict, slide_num: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, colors["bg"])

        # Slide number indicator
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(1), Inches(0.4))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"{slide_num}"
        p.font.size = Pt(14)
        p.font.color.rgb = colors["primary"]
        p.font.bold = True

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.LEFT

        # Underline
        self._add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), colors["secondary"])

        # Content
        bullets = slide_data.get("bullets", [])
        if bullets:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf2.paragraphs[0]
                else:
                    p = tf2.add_paragraph()
                p.text = f"▸  {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = colors["text"]
                p.space_after = Pt(12)
                p.level = 0

        # Accent dots on right side
        for i in range(3):
            y = Inches(2.0 + i * 1.5)
            self._add_circle(slide, Inches(12.0), y, Inches(0.15), colors["primary"])

    def _add_closing_slide(self, prs, colors: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, colors["bg"])

        # Accent line
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), colors["primary"])

        # Thank you
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Danke für Ihre Aufmerksamkeit!"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = "Fragen? | Erstellt mit J.A.R.V.I.S."
        p2.font.size = Pt(18)
        p2.font.color.rgb = colors["text_dim"]
        p2.alignment = PP_ALIGN.CENTER

        # Bottom accent line
        self._add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), colors["primary"])

    def _set_slide_bg(self, slide, color: RGBColor):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_shape(self, slide, left, top, width, height, color: RGBColor):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_circle(self, slide, left, top, size, color: RGBColor):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape


# Convenience function
def create_presentation(topic: str, slides: int = 8, theme: str = "jarvis") -> str:
    mgr = PresentationManager()
    return mgr.create_presentation(topic, slides, theme)